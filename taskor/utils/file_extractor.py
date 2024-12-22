import json
import os
import re
import warnings
import xml.etree.ElementTree as ET
import zipfile
from configparser import ConfigParser
import ebooklib
import PyPDF2
import markdown2
import openpyxl
import pandas as pd
import pytesseract
import xlrd
import yaml
from bs4 import BeautifulSoup
from docx import Document
from ebooklib import epub
from PIL import Image
from pptx import Presentation
import striprtf


warnings.filterwarnings("ignore", category=UserWarning, module="ebooklib")
warnings.filterwarnings("ignore", category=FutureWarning, module="ebooklib")


def extract_text_from_file(file_path):
    file_ext = os.path.splitext(file_path)[1].lower()
    text = ""

    try:
        if file_ext == ".txt":
            text = extract_text_from_txt(file_path)
        elif file_ext == ".pdf":
            text = extract_text_from_pdf(file_path)
        elif file_ext in (".doc", ".docx"):
            text = extract_text_from_doc(file_path)
        elif file_ext in (".xls", ".xlsx", ".csv", ".ods"):
            text = extract_text_from_spreadsheet(file_path)
        elif file_ext in (".ppt", ".pptx", ".odp"):
            text = extract_text_from_ppt(file_path)
        elif file_ext in (".html", ".htm"):
            text = extract_text_from_html(file_path)
        elif file_ext in (".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".gif"):
            text = extract_text_from_image(file_path)
        elif file_ext in (".epub", ".mobi", ".azw3"):
            text = extract_text_from_epub(file_path)
        elif file_ext == ".rtf":
            text = extract_text_from_rtf(file_path)
        elif file_ext == ".odt":
            text = extract_text_from_odt(file_path)
        elif file_ext == ".md":
            text = extract_text_from_markdown(file_path)
        elif file_ext == ".json":
            text = extract_text_from_json(file_path)
        elif file_ext in (".yaml", ".yml"):
            text = extract_text_from_yaml(file_path)
        elif file_ext == ".xml":
            text = extract_text_from_xml(file_path)
        elif file_ext == ".ini":
            text = extract_text_from_ini(file_path)
    except:
        pass

    if text == "":
        try:
            text = extract_text_from_binary(file_path)
        except:
            pass

    return file_path + "\n" + text


def extract_text_from_txt(file_path: str) -> str:

    with open(file_path, encoding="utf-8", errors="ignore") as f:
        return f.read()


def extract_text_from_pdf(file_path: str) -> str:
    text = ""
    with open(file_path, "rb") as file:
        reader = PyPDF2.PdfReader(file)
        for page_num, page in enumerate(reader.pages, start=1):
            page_text = page.extract_text()
            if page_text:
                text += f"\n--- Page {page_num} ---\n{page_text}"

    return text


def extract_text_from_doc(file_path: str) -> str:
    doc = Document(file_path)
    text = "\n".join([para.text for para in doc.paragraphs])

    return text


def extract_text_from_spreadsheet(file_path: str) -> str:
    text = ""

    if file_path.endswith((".xlsx", ".xlsm")):
        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        for sheetname in wb.sheetnames:
            ws = wb[sheetname]
            text += f"\n--- Sheet: {sheetname} ---\n"
            for row in ws.iter_rows(values_only=True):
                row_text = " | ".join(
                    [str(cell) if cell is not None else "" for cell in row]
                )
                text += row_text + "\n"
    elif file_path.endswith(".xls"):
        wb = xlrd.open_workbook(file_path)
        for sheet in wb.sheets():
            text += f"\n--- Sheet: {sheet.name} ---\n"
            for row_idx in range(sheet.nrows):
                row = sheet.row(row_idx)
                row_text = " | ".join(
                    [
                        str(cell.value) if cell.value is not None else ""
                        for cell in row
                    ]
                )
                text += row_text + "\n"
    elif file_path.endswith(".csv"):
        df = pd.read_csv(file_path)
        text += df.to_string(index=False)
    elif file_path.endswith(".ods"):
        df = pd.read_excel(file_path, engine="odf")
        text += df.to_string(index=False)

    return text


def extract_text_from_ppt(file_path: str) -> str:
    text_runs = []
    prs = Presentation(file_path)
    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_text = f"\n--- Slide {slide_num} ---\n"
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text += shape.text + "\n"
        text_runs.append(slide_text)

    return "\n".join(text_runs)


def extract_text_from_html(file_path: str) -> str:
    with open(file_path, encoding="utf-8", errors="ignore") as f:
        html_content = f.read()
    soup = BeautifulSoup(html_content, "html.parser")
    text = soup.get_text(separator="\n")

    return text


def extract_text_from_image(file_path: str) -> str:
    img = Image.open(file_path)
    text = pytesseract.image_to_string(img)

    return text


def extract_text_from_epub(file_path: str) -> str:
    book = epub.read_epub(file_path)
    text = ""

    for item in book.get_items():
        if item.get_type() == ebooklib.ITEM_DOCUMENT:
            content = item.get_content()
            soup = BeautifulSoup(content, "html.parser")
            text += soup.get_text(separator="\n")
    return text



def extract_text_from_rtf(file_path: str) -> str:
    with open(file_path, encoding="utf-8", errors="ignore") as f:
        rtf_content = f.read()
    text = striprtf.rtf_to_text(rtf_content)

    return text

def extract_text_from_odt(file_path: str) -> str:
    with zipfile.ZipFile(file_path) as z:
        with z.open("content.xml") as content_file:
            tree = ET.parse(content_file)
            root = tree.getroot()
            namespaces = {
                "text": "urn:oasis:names:tc:opendocument:xmlns:text:1.0",
            }
            texts = root.findall(".//text:p", namespaces)
            text = "\n".join([t.text for t in texts if t.text])

    return text

def extract_text_from_markdown(file_path: str) -> str:
    with open(file_path, encoding="utf-8", errors="ignore") as f:
        md_content = f.read()
    html = markdown2.markdown(md_content)
    soup = BeautifulSoup(html, "html.parser")
    text = soup.get_text(separator="\n")

    return text


def extract_text_from_json(file_path: str) -> str:
    with open(file_path, encoding="utf-8", errors="ignore") as f:
        data = json.load(f)
    text = json.dumps(data, indent=2)

    return text


def extract_text_from_yaml(file_path: str) -> str:
    with open(file_path, encoding="utf-8", errors="ignore") as f:
        data = yaml.safe_load(f)
    text = yaml.dump(data)

    return text

def extract_text_from_xml(file_path: str) -> str:
    tree = ET.parse(file_path)
    root = tree.getroot()
    text_elements = []

    def recursive_text_extract(element):
        if element.text and element.text.strip():
            text_elements.append(element.text.strip())
        for child in element:
            recursive_text_extract(child)

    recursive_text_extract(root)

    return "\n".join(text_elements)

def extract_text_from_ini(file_path: str) -> str:
    parser = ConfigParser()
    parser.read(file_path)
    text = ""

    for section in parser.sections():
        text += f"[{section}]\n"
        for key, value in parser.items(section):
            text += f"{key} = {value}\n"
        text += "\n"

    return text


def extract_text_from_binary(file_path: str) -> str:
    with open(file_path, "rb") as f:
        content = f.read()

    printable_text = re.findall(b"[\x20-\x7E]+", content)

    decoded_text = "\n".join(
        [txt.decode("utf-8", errors="replace") for txt in printable_text]
    )

    return decoded_text

if __name__ == "__main__":
    print(extract_text_from_file("../taskor.py"))