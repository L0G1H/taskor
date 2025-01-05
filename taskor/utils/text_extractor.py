from __future__ import annotations
import json
import re
import warnings
import xml.etree.ElementTree as ET
import zipfile
from configparser import ConfigParser
import ebooklib
import PyPDF2
import markdown2
import openpyxl
import pytesseract
import yaml
from bs4 import BeautifulSoup
from docx import Document
from ebooklib import epub
from PIL import Image
from pptx import Presentation
import striprtf
from pathlib import Path


warnings.filterwarnings("ignore", category=UserWarning, module="ebooklib")
warnings.filterwarnings("ignore", category=FutureWarning, module="ebooklib")


def extract_text_from_file(file_path: str) -> str | tuple[None, str]:
    file_ext = Path(file_path).suffix.lower()

    try:
        if file_ext == ".txt":
            text = extract_text_from_txt(file_path)
        elif file_ext == ".pdf":
            text = extract_text_from_pdf(file_path)
        elif file_ext in (".doc", ".docx"):
            text = extract_text_from_doc(file_path)
        elif file_ext in (".xls", ".xlsx", ".csv", ".ods"):
            text = extract_text_from_spreadsheet(file_path)
        elif file_ext in (".ppt", ".pptx", ".odp"):
            text = extract_text_from_ppt(file_path)
        elif file_ext in (".html", ".htm"):
            text = extract_text_from_html(file_path)
        elif file_ext in (".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".gif"):
            text = extract_text_from_image(file_path)
        elif file_ext in (".epub", ".mobi", ".azw3"):
            text = extract_text_from_epub(file_path)
        elif file_ext == ".rtf":
            text = extract_text_from_rtf(file_path)
        elif file_ext == ".odt":
            text = extract_text_from_odt(file_path)
        elif file_ext == ".md":
            text = extract_text_from_markdown(file_path)
        elif file_ext == ".json":
            text = extract_text_from_json(file_path)
        elif file_ext in (".yaml", ".yml"):
            text = extract_text_from_yaml(file_path)
        elif file_ext == ".xml":
            text = extract_text_from_xml(file_path)
        elif file_ext == ".ini":
            text = extract_text_from_ini(file_path)
        else:
            text = extract_text_from_binary(file_path)
    except Exception as e:
        return None, str(e)

    return file_path + "\n" + text


def extract_text_from_txt(file_path: str) -> str:
    with Path.open(file_path, encoding="utf-8", errors="ignore") as f:
        return f.read()


def extract_text_from_pdf(file_path: str) -> str:
    text = []
    with Path.open(file_path, "rb") as file:
        reader = PyPDF2.PdfReader(file)
        for page_num, page in enumerate(reader.pages, start=1):
            try:
                page_text = page.extract_text()
                if page_text and not page_text.isspace():
                    text.append(f"\n--- Page {page_num} ---\n{page_text.strip()}")
            except Exception as e:
                text.append(f"\n--- Page {page_num} ---\nError: {e}")

    return "\n".join(text)


def extract_text_from_doc(file_path: str) -> str:
    doc = Document(file_path)

    return "\n".join([para.text for para in doc.paragraphs])


def extract_text_from_spreadsheet(file_path: str) -> str:
    def _format_cell(cell: None | float | str) -> str:
        if cell is None:
            return ""
        if isinstance(cell, (int, float)):
            return str(cell)
        return str(cell).strip()

    text = []

    if file_path.endswith((".xlsx", ".xlsm")):
        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        for sheetname in wb.sheetnames:
            ws = wb[sheetname]
            text.append(f"\n--- Sheet: {sheetname} ---")
            for row in ws.iter_rows(values_only=True):
                if any(cell is not None for cell in row):
                    row_text = " | ".join(_format_cell(cell) for cell in row)
                    text.append(row_text)

    return "\n".join(text)


def extract_text_from_ppt(file_path: str) -> str:
    text_runs = []
    prs = Presentation(file_path)
    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_text = f"\n--- Slide {slide_num} ---\n"
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text += shape.text + "\n"
        text_runs.append(slide_text)

    return "\n".join(text_runs)


def extract_text_from_html(file_path: str) -> str:
    with Path.open(file_path, encoding="utf-8", errors="ignore") as f:
        html_content = f.read()
    soup = BeautifulSoup(html_content, "html.parser")

    return soup.get_text(separator="\n")


def extract_text_from_image(file_path: str) -> str:
    img = Image.open(file_path)

    return pytesseract.image_to_string(img)


def extract_text_from_epub(file_path: str) -> str:
    book = epub.read_epub(file_path)
    text = ""

    for item in book.get_items():
        if item.get_type() == ebooklib.ITEM_DOCUMENT:
            content = item.get_content()
            soup = BeautifulSoup(content, "html.parser")
            text += soup.get_text(separator="\n")

    return text


def extract_text_from_rtf(file_path: str) -> str:
    with Path.open(file_path, encoding="utf-8", errors="ignore") as f:
        rtf_content = f.read()

    return striprtf.rtf_to_text(rtf_content)


def extract_text_from_odt(file_path: str) -> str:
    with zipfile.ZipFile(file_path) as z, z.open("content.xml") as content_file:
        tree = ET.parse(content_file)
        root = tree.getroot()
        namespaces = {
            "text": "urn:oasis:names:tc:opendocument:xmlns:text:1.0",
        }
        texts = root.findall(".//text:p", namespaces)

    return "\n".join([t.text for t in texts if t.text])


def extract_text_from_markdown(file_path: str) -> str:
    with Path.open(file_path, encoding="utf-8", errors="ignore") as f:
        md_content = f.read()

    html = markdown2.markdown(md_content)
    soup = BeautifulSoup(html, "html.parser")

    return soup.get_text(separator="\n")


def extract_text_from_json(file_path: str) -> str:
    with Path.open(file_path, encoding="utf-8", errors="ignore") as f:
        data = json.load(f)

    return json.dumps(data, indent=2)


def extract_text_from_yaml(file_path: str) -> str:
    with Path.open(file_path, encoding="utf-8", errors="ignore") as f:
        data = yaml.safe_load(f)

    return yaml.dump(data)


def extract_text_from_xml(file_path: str) -> str:
    tree = ET.parse(file_path)
    root = tree.getroot()
    text_elements = []

    def recursive_text_extract(element: ET.Element) -> None:
        if element.text and element.text.strip():
            text_elements.append(element.text.strip())
        for child in element:
            recursive_text_extract(child)

    recursive_text_extract(root)

    return "\n".join(text_elements)


def extract_text_from_ini(file_path: str) -> str:
    parser = ConfigParser()
    parser.read(file_path)
    text = ""

    for section in parser.sections():
        text += f"[{section}]\n"
        for key, value in parser.items(section):
            text += f"{key} = {value}\n"
        text += "\n"

    return text


def extract_text_from_binary(file_path: str) -> str:
    encodings = ["utf-8", "latin-1", "ascii"]
    min_length = 4

    with Path.open(file_path, "rb") as f:
        content = f.read()

    text_pattern = re.compile(b"[\x20-\x7e]{%d,}" % min_length)
    printable_text = text_pattern.findall(content)

    results = []
    for text in printable_text:
        for encoding in encodings:
            try:
                decoded = text.decode(encoding)
                if not decoded.isspace():
                    results.append(decoded)
                break
            except UnicodeError:
                continue

    return "\n".join(results)


if __name__ == "__main__":
    print(extract_text_from_file("../taskor.py"))
